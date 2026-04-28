"""
TEM Auto-Processor — Rotate, Crop & PPT Generator
====================================================
Process TEM images: auto-rotate (FFT), crop, and generate PPT.

Images are placed in two folders to indicate type:
    standard/   — Standard (cross-section) TEM images
    planar/     — Planar TEM images

Usage:
    python tem_process.py <input_folder>          # auto-detect standard/ and planar/ subfolders
    python tem_process.py                          # use current directory
    python tem_process.py <folder> -o result.pptx  # custom output path

Requirements:
    pip install python-pptx Pillow opencv-python-headless numpy
"""
import os
import sys
import math
import glob
import io
import argparse

import numpy as np
import cv2
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


# ─── Configuration ────────────────────────────────────────
SCALE_BAR_CROP_RATIO = 0.06

PPT_SPECS = {
    "standard": {
        "resize": (3.8, 3.8),
        "crop": (2.8, 2.5),
    },
    "planar": {
        "resize": (2.8, 2.8),
        "crop": (1.81, 2.5),
    },
}


# ═══════════════════════════════════════════════════════════
# Image Processing
# ═══════════════════════════════════════════════════════════

def load_tem_image(path):
    """Load TIF, convert to grayscale, crop scale bar."""
    img = Image.open(path)
    if img.mode != "L":
        img = img.convert("L")
    w, h = img.size
    crop_h = int(h * (1 - SCALE_BAR_CROP_RATIO))
    img = img.crop((0, 0, w, crop_h))
    return img


def detect_rotation_fft(img_pil, image_type):
    """Detect rotation angle via FFT power spectrum analysis.

    Standard TEM: horizontal layers -> FFT peak near 90 deg.
    Planar TEM:   vertical structures -> FFT peak near 0 deg.

    Returns: correction angle in degrees.
    """
    gray = np.array(img_pil, dtype=np.float32)

    # Downscale to 512px for speed
    scale = 512.0 / max(gray.shape)
    small = cv2.resize(gray, None, fx=scale, fy=scale, interpolation=cv2.INTER_AREA)
    h, w = small.shape

    # Hanning window to suppress spectral leakage at edges
    win_r = np.hanning(h).reshape(-1, 1)
    win_c = np.hanning(w).reshape(1, -1)
    windowed = small * win_r * win_c

    # 2D FFT -> power spectrum
    fshift = np.fft.fftshift(np.fft.fft2(windowed))
    magnitude = np.log(np.abs(fshift) + 1.0)

    cy, cx = h // 2, w // 2

    # Band-pass: keep only mid-frequency (skip DC and high-freq noise)
    Y, X = np.ogrid[:h, :w]
    r = np.sqrt((X - cx) ** 2 + (Y - cy) ** 2)
    magnitude *= (r > 5) & (r < min(h, w) // 3)

    # Sweep angles: look for peak power direction
    if image_type == "standard":
        search = np.arange(70.0, 110.0, 0.1)
        ref = 90.0
    else:
        search = np.arange(-20.0, 20.0, 0.1)
        ref = 0.0

    rs = np.arange(5, min(h, w) // 3)
    best_angle = ref
    best_power = 0.0

    for adeg in search:
        arad = math.radians(adeg)
        xs = np.clip((cx + rs * math.cos(arad)).astype(int), 0, w - 1)
        ys = np.clip((cy + rs * math.sin(arad)).astype(int), 0, h - 1)
        power = float(np.sum(magnitude[ys, xs]))
        if power > best_power:
            best_power = power
            best_angle = adeg

    if image_type == "standard":
        return -(best_angle - 90.0)
    else:
        return -best_angle


def rotate_and_crop(img, angle_deg, spec):
    """Rotate image, remove black borders, resize and crop to PPT spec."""
    rotated = img.rotate(-angle_deg, resample=Image.BICUBIC, expand=True, fillcolor=0)

    w, h = rotated.size
    rad = abs(math.radians(angle_deg))
    orig_w, orig_h = img.size

    if rad < 1e-6:
        crop_w, crop_h = orig_w, orig_h
    else:
        cos_a, sin_a = abs(math.cos(rad)), abs(math.sin(rad))
        if orig_w <= 2.0 * orig_h * sin_a * cos_a or orig_h <= 2.0 * orig_w * sin_a * cos_a:
            crop_w = int(orig_h / (2 * sin_a))
            crop_h = int(orig_w / (2 * sin_a))
        else:
            crop_w = int((orig_w * cos_a - orig_h * sin_a) / (cos_a**2 - sin_a**2))
            crop_h = int((orig_h * cos_a - orig_w * sin_a) / (cos_a**2 - sin_a**2))
        crop_w = min(crop_w, w)
        crop_h = min(crop_h, h)

    left = (w - crop_w) // 2
    top = (h - crop_h) // 2
    rotated = rotated.crop((left, top, left + crop_w, top + crop_h))

    dpi = 300
    rw, rh = int(spec["resize"][0] * dpi), int(spec["resize"][1] * dpi)
    resized = rotated.resize((rw, rh), Image.LANCZOS)

    cw, ch = int(spec["crop"][0] * dpi), int(spec["crop"][1] * dpi)
    left = (rw - cw) // 2
    top = (rh - ch) // 2
    return resized.crop((left, top, left + cw, top + ch))


# ═══════════════════════════════════════════════════════════
# PPT Generation
# ═══════════════════════════════════════════════════════════

def generate_ppt(results, output_path):
    """Generate PPT with processed TEM images, grouped by type, multiple per slide."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    grouped = {"standard": [], "planar": []}
    for r in results:
        grouped[r["type"]].append(r)

    for type_name, items in grouped.items():
        if not items:
            continue

        spec = PPT_SPECS[type_name]
        img_w = spec["crop"][0]
        img_h = spec["crop"][1]
        label_h = 0.3

        margin_x = 0.3
        margin_top = 0.8
        gap_x = 0.25
        gap_y = 0.15
        usable_w = 13.333 - 2 * margin_x
        usable_h = 7.5 - margin_top - 0.2

        cols = max(1, int((usable_w + gap_x) / (img_w + gap_x)))
        rows = max(1, int((usable_h + gap_y) / (img_h + label_h + gap_y)))
        per_slide = cols * rows

        grid_w = cols * img_w + (cols - 1) * gap_x
        grid_h = rows * (img_h + label_h) + (rows - 1) * gap_y
        start_x = (13.333 - grid_w) / 2
        start_y = margin_top + (usable_h - grid_h) / 2

        for page_idx in range(0, len(items), per_slide):
            page_items = items[page_idx:page_idx + per_slide]
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            txBox = slide.shapes.add_textbox(
                Inches(0.3), Inches(0.15), Inches(12), Inches(0.5))
            p = txBox.text_frame.paragraphs[0]
            page_num = page_idx // per_slide + 1
            total_pages = (len(items) + per_slide - 1) // per_slide
            p.text = f"{type_name.title()} TEM ({len(items)} images) \u2014 Page {page_num}/{total_pages}"
            p.font.size = Pt(16)
            p.font.bold = True

            for i, r in enumerate(page_items):
                row = i // cols
                col = i % cols
                x = start_x + col * (img_w + gap_x)
                y = start_y + row * (img_h + label_h + gap_y)

                buf = io.BytesIO()
                r["processed_image"].save(buf, format="PNG")
                buf.seek(0)
                slide.shapes.add_picture(
                    buf, Inches(x), Inches(y), Inches(img_w), Inches(img_h))

                txBox2 = slide.shapes.add_textbox(
                    Inches(x), Inches(y + img_h), Inches(img_w), Inches(label_h))
                p2 = txBox2.text_frame.paragraphs[0]
                p2.text = r["filename"]
                p2.font.size = Pt(12)
                p2.alignment = 1  # center

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    for type_name in ["standard", "planar"]:
        n = len(grouped[type_name])
        if n > 0:
            print(f"  {type_name.title()}: {n} images")
    print(f"\nPPT saved: {output_path}")


# ═══════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════

def collect_tifs(folder):
    """Collect all TIF/TIFF files from a folder."""
    if not folder or not os.path.isdir(folder):
        return []
    return sorted(
        glob.glob(os.path.join(folder, "*.tif")) +
        glob.glob(os.path.join(folder, "*.tiff"))
    )


def process_images(typed_paths):
    """Process TEM images. typed_paths = list of (path, type_name)."""
    results = []
    total = len(typed_paths)
    for i, (path, type_name) in enumerate(typed_paths):
        try:
            img = load_tem_image(path)
            angle = detect_rotation_fft(img, type_name)
            processed = rotate_and_crop(img, angle, PPT_SPECS[type_name])

            results.append({
                "filename": os.path.basename(path),
                "type": type_name,
                "angle": angle,
                "processed_image": processed,
            })
            print(f"  [{i+1}/{total}] {os.path.basename(path)}: "
                  f"{type_name}, angle={angle:.2f}")

        except Exception as e:
            print(f"  [{i+1}/{total}] ERROR {os.path.basename(path)}: {e}")

    return results


def main():
    parser = argparse.ArgumentParser(
        description="TEM Auto-Processor: auto-rotate, crop, and generate PPT"
    )
    parser.add_argument("input", nargs="?", default=".",
                        help="Input folder containing standard/ and planar/ subfolders (default: current dir)")
    parser.add_argument("-o", "--output", default=None,
                        help="Output PPT path (default: <input>/output.pptx)")
    args = parser.parse_args()

    root = os.path.abspath(args.input)
    if not os.path.isdir(root):
        print(f"ERROR: {root} is not a directory.")
        sys.exit(1)

    # Only match exact folder names: standard/ and planar/
    std_dir = pln_dir = None
    for name in os.listdir(root):
        full = os.path.join(root, name)
        if not os.path.isdir(full):
            continue
        low = name.lower()
        if low == "standard":
            std_dir = full
        elif low == "planar":
            pln_dir = full

    if not std_dir and not pln_dir:
        print(f"ERROR: No 'standard' or 'planar' subfolder found in {root}")
        print("Expected structure:")
        print(f"  {root}/")
        print(f"    standard/   <- Standard TEM TIF files")
        print(f"    planar/     <- Planar TEM TIF files")
        sys.exit(1)

    # Collect images with their type
    typed_paths = []
    for folder, type_name in [(std_dir, "standard"), (pln_dir, "planar")]:
        tifs = collect_tifs(folder)
        if tifs:
            print(f"Found {len(tifs)} {type_name} TIF files in: {folder}")
        typed_paths.extend((p, type_name) for p in tifs)

    if not typed_paths:
        print("ERROR: No TIF files found in the provided folders.")
        sys.exit(1)

    print(f"\nProcessing {len(typed_paths)} images...")
    results = process_images(typed_paths)

    n_std = sum(1 for r in results if r["type"] == "standard")
    n_pln = sum(1 for r in results if r["type"] == "planar")
    print(f"\nProcessed: {n_std} Standard, {n_pln} Planar")

    output_path = args.output
    if not output_path:
        output_path = os.path.join(root, "output.pptx")
    elif not os.path.isabs(output_path):
        output_path = os.path.abspath(output_path)
    generate_ppt(results, output_path)


if __name__ == "__main__":
    main()
